from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1a, 0x2a, 0x4a)
ORANGE = RGBColor(0xe8, 0x9b, 0x3e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

# ============ SLIDE 1: Title Slide ============
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK_BLUE)

# Orange accent bar at top
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ORANGE)

# Company name
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "KP Travels", font_size=54, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "Bring the Happiness", font_size=32, bold=False, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(2), Inches(4.0), Inches(9), Inches(0.6),
             "Your Trusted Partner for Safe & Comfortable Travel in Bangalore", font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

# Orange line divider
add_shape_bg(slide, Inches(5.5), Inches(3.85), Inches(2.3), Inches(0.04), ORANGE)

# Contact info
add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.5),
             "Phone: +91 99454 98275  |  Email: admin@actoursandtravels.com", font_size=16, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

# Website
add_text_box(slide, Inches(2), Inches(6.0), Inches(9), Inches(0.5),
             "www.actoursandtravels.com", font_size=16, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Orange accent bar at bottom
add_shape_bg(slide, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), ORANGE)


# ============ SLIDE 2: About Us ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Header bar
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "About KP Travels", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

# Orange underline
add_shape_bg(slide, Inches(0.5), Inches(1.2), Inches(3), Inches(0.06), ORANGE)

# About text
about_text = ("KP Travels Solutions Pvt. Ltd. is a trusted transportation service provider based in Bangalore. "
              "We specialize in providing safe, comfortable, and reliable travel solutions for "
              "individuals and corporate clients.\n\n"
              "Our team of professionally trained drivers ensures that every journey is smooth and pleasant. "
              "We maintain a fleet of well-serviced vehicles to cater to all your travel needs.")
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(7), Inches(2.5),
             about_text, font_size=17, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

# Key highlights box
highlight_box = add_shape_bg(slide, Inches(8.5), Inches(1.6), Inches(4.2), Inches(4.5), RGBColor(0xF0, 0xF4, 0xF8))
highlight_box.shadow.inherit = False

add_text_box(slide, Inches(8.8), Inches(1.8), Inches(3.8), Inches(0.5),
             "Key Highlights", font_size=22, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)

highlights = [
    "✔  Trained Professional Drivers",
    "✔  Well-Maintained Vehicles",
    "✔  Transparent Pricing",
    "✔  On-time Service",
    "✔  GPS Enabled Vehicles",
    "✔  24/7 Customer Support"
]
for i, h in enumerate(highlights):
    add_text_box(slide, Inches(8.8), Inches(2.5 + i * 0.5), Inches(3.8), Inches(0.5),
                 h, font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

# Stats row
stats = [("50+", "Professional\nDrivers"), ("1000+", "Happy\nCustomers"), ("30+", "Vehicles"), ("5+", "Years\nExperience")]
for i, (num, label) in enumerate(stats):
    x = Inches(0.8 + i * 1.8)
    stat_box = add_shape_bg(slide, x, Inches(5.0), Inches(1.5), Inches(1.8), DARK_BLUE)
    add_text_box(slide, x, Inches(5.2), Inches(1.5), Inches(0.7),
                 num, font_size=32, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(5.9), Inches(1.5), Inches(0.7),
                 label, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)


# ============ SLIDE 3: Our Services ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "Our Services", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
add_shape_bg(slide, Inches(0.5), Inches(1.2), Inches(2.5), Inches(0.06), ORANGE)

services = [
    ("Local Rentals", "Hourly and daily car rentals for city travel. Perfect for meetings, shopping, or personal errands within Bangalore."),
    ("Airport Transfer", "Punctual pickup and drop services to Kempegowda International Airport. Track your flight, we'll be there on time."),
    ("Outstation Trips", "Travel to Mysore, Coorg, Ooty, Goa and more. One-way and round-trip packages at competitive prices."),
    ("Corporate Travel", "Employee transportation for IT companies. Dedicated fleet, trained chauffeurs, and monthly billing options."),
    ("Wedding & Events", "Make your special day memorable. Decorated cars for weddings, parties, and special occasions."),
    ("24/7 Support", "Book anytime, travel anytime. Our support team is available round the clock for your convenience."),
]

icons = ["🕐", "✈️", "🏠", "👥", "❤️", "📞"]

for i, (title, desc) in enumerate(services):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.6 + row * 2.7)

    card = add_shape_bg(slide, x, y, Inches(3.8), Inches(2.3), RGBColor(0xF8, 0xF9, 0xFA))

    # Icon circle
    icon_circle = add_shape_bg(slide, x + Inches(0.2), y + Inches(0.25), Inches(0.6), Inches(0.6), ORANGE)

    add_text_box(slide, x + Inches(1.0), y + Inches(0.3), Inches(2.6), Inches(0.4),
                 title, font_size=18, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.4), Inches(1.2),
                 desc, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)


# ============ SLIDE 4: Our Fleet ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "Our Fleet", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
add_shape_bg(slide, Inches(0.5), Inches(1.2), Inches(2), Inches(0.06), ORANGE)

# 4 Seaters
add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(0.5),
             "4 Seaters", font_size=24, bold=True, color=ORANGE, alignment=PP_ALIGN.LEFT)

four_seaters = [
    ("Maruti Swift / Dzire", "4 Seater | AC | Comfortable", "Best for City Travel"),
    ("Toyota Etios", "4 Seater | AC | Reliable", "Best for Airport & Outstation"),
]

for i, (name, spec, best) in enumerate(four_seaters):
    x = Inches(0.5 + i * 3.5)
    card = add_shape_bg(slide, x, Inches(2.1), Inches(3.2), Inches(1.8), RGBColor(0xF0, 0xF4, 0xF8))
    add_text_box(slide, x + Inches(0.2), Inches(2.3), Inches(2.8), Inches(0.4),
                 name, font_size=18, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, x + Inches(0.2), Inches(2.8), Inches(2.8), Inches(0.4),
                 spec, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, x + Inches(0.2), Inches(3.3), Inches(2.8), Inches(0.4),
                 best, font_size=14, bold=True, color=ORANGE, alignment=PP_ALIGN.LEFT)

# 7+ Seaters
add_text_box(slide, Inches(0.5), Inches(4.2), Inches(4), Inches(0.5),
             "7+ Seaters", font_size=24, bold=True, color=ORANGE, alignment=PP_ALIGN.LEFT)

seven_seaters = [
    ("Maruti Ertiga", "6-7 Seater | AC | Family Friendly", "Best for Family Trips"),
    ("Toyota Innova Crysta", "6-7 Seater | AC | Spacious", "Best for Groups & Outstation"),
    ("Tempo Traveller", "12-17 Seater | AC | Spacious", "Best for Large Groups & Tours"),
]

for i, (name, spec, best) in enumerate(seven_seaters):
    x = Inches(0.5 + i * 4.2)
    card = add_shape_bg(slide, x, Inches(4.8), Inches(3.8), Inches(1.8), RGBColor(0xF0, 0xF4, 0xF8))
    add_text_box(slide, x + Inches(0.2), Inches(5.0), Inches(3.4), Inches(0.4),
                 name, font_size=18, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, x + Inches(0.2), Inches(5.5), Inches(3.4), Inches(0.4),
                 spec, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, x + Inches(0.2), Inches(6.0), Inches(3.4), Inches(0.4),
                 best, font_size=14, bold=True, color=ORANGE, alignment=PP_ALIGN.LEFT)


# ============ SLIDE 5: Why Choose Us ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "Why Choose Us", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
add_shape_bg(slide, Inches(0.5), Inches(1.2), Inches(2.8), Inches(0.06), ORANGE)

reasons = [
    ("Punctuality", "We value your time. Our drivers arrive before scheduled pickup time."),
    ("Safety First", "Verified drivers, sanitized vehicles, and GPS tracking for every trip."),
    ("Best Prices", "Transparent pricing with no hidden charges. Competitive rates guaranteed."),
    ("24/7 Service", "Book anytime, travel anytime. We're available round the clock."),
]

for i, (title, desc) in enumerate(reasons):
    x = Inches(0.5 + i * 3.2)

    # Card background
    card = add_shape_bg(slide, x, Inches(1.8), Inches(2.9), Inches(3.0), DARK_BLUE)

    # Orange circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.95), Inches(2.1), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()

    add_text_box(slide, x, Inches(3.3), Inches(2.9), Inches(0.5),
                 title, font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.9), Inches(2.5), Inches(0.8),
                 desc, font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

# Bottom tagline
add_shape_bg(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(1.2), RGBColor(0xF8, 0xF9, 0xFA))
add_text_box(slide, Inches(1), Inches(5.7), Inches(11.3), Inches(0.8),
             "\"We don't just provide rides, we deliver experiences. Your comfort and safety are our top priority.\"",
             font_size=18, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)


# ============ SLIDE 6: Popular Destinations ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "Popular Destinations", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
add_shape_bg(slide, Inches(0.5), Inches(1.2), Inches(3.5), Inches(0.06), ORANGE)

destinations = [
    ("Mysore", "Heritage palaces, royal city\n~150 km from Bangalore"),
    ("Coorg", "Coffee plantations, waterfalls\n~260 km from Bangalore"),
    ("Ooty", "Hill station, scenic beauty\n~270 km from Bangalore"),
    ("Goa", "Beaches, nightlife, churches\n~560 km from Bangalore"),
    ("Tirupati", "Temple town, spiritual hub\n~250 km from Bangalore"),
    ("Hampi", "UNESCO heritage, ruins\n~340 km from Bangalore"),
]

for i, (place, desc) in enumerate(destinations):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.6 + row * 2.7)

    card = add_shape_bg(slide, x, y, Inches(3.8), Inches(2.3), DARK_BLUE)

    add_text_box(slide, x, y + Inches(0.4), Inches(3.8), Inches(0.6),
                 place, font_size=28, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.1), Inches(3.8), Inches(1.0),
                 desc, font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER)


# ============ SLIDE 7: Our Trusted Clients ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "Our Trusted Clients", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
add_shape_bg(slide, Inches(0.5), Inches(1.2), Inches(3.5), Inches(0.06), ORANGE)

add_text_box(slide, Inches(1), Inches(1.6), Inches(11), Inches(0.6),
             "Proud to serve leading companies and organizations",
             font_size=20, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

clients = [
    ("Pluto Rides", "Ride-hailing & mobility platform"),
    ("Diligent", "Global governance & compliance leader"),
    ("OpenText", "Enterprise information management"),
    ("Operations Insights", "Business analytics & operations"),
]

for i, (name, desc) in enumerate(clients):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.8)

    # Client card
    card = add_shape_bg(slide, x, y, Inches(2.8), Inches(3.0), RGBColor(0xF0, 0xF4, 0xF8))

    # Orange accent at top of card
    add_shape_bg(slide, x, y, Inches(2.8), Inches(0.06), ORANGE)

    # Client name
    add_text_box(slide, x, y + Inches(0.8), Inches(2.8), Inches(0.6),
                 name, font_size=22, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    # Orange divider
    add_shape_bg(slide, x + Inches(0.9), y + Inches(1.5), Inches(1.0), Inches(0.04), ORANGE)

    # Description
    add_text_box(slide, x + Inches(0.2), y + Inches(1.8), Inches(2.4), Inches(0.8),
                 desc, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom note
add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.6),
             "We provide dedicated corporate transportation solutions tailored to each client's needs.",
             font_size=16, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)


# ============ SLIDE 8: Contact Us ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ORANGE)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
             "Let's Work Together", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(5.5), Inches(2.1), Inches(2.3), Inches(0.04), ORANGE)

add_text_box(slide, Inches(2), Inches(2.5), Inches(9), Inches(0.6),
             "Contact us for corporate tie-ups, bulk bookings, and vendor partnerships",
             font_size=18, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

# Contact cards
contact_items = [
    ("Phone", "+91 99454 98275"),
    ("Email", "admin@actoursandtravels.com"),
    ("Website", "www.actoursandtravels.com"),
    ("WhatsApp", "+91 99454 98275"),
]

for i, (label, value) in enumerate(contact_items):
    x = Inches(1.5 + i * 2.8)
    card = add_shape_bg(slide, x, Inches(3.5), Inches(2.5), Inches(1.8), RGBColor(0x24, 0x3a, 0x5e))

    add_text_box(slide, x, Inches(3.7), Inches(2.5), Inches(0.5),
                 label, font_size=16, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(4.3), Inches(2.5), Inches(0.8),
                 value, font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER)

# Address
add_text_box(slide, Inches(2), Inches(5.8), Inches(9), Inches(0.5),
             "Address: No.33-2, Kudlu Main Road, Kudlu, Bommanahalli Post, Bangalore - 560068",
             font_size=16, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

# Bottom bar
add_shape_bg(slide, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), ORANGE)

# Save
output_path = "/Users/govindpawar/Plt/ac-travels/KP_Travels_Vendor_Presentation.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
